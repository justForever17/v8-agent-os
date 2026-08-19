import os
from pathlib import Path
import json
import csv
import zipfile
import tarfile
import tempfile
import logging

try:
    import pymupdf
except ImportError:
    pymupdf = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import xlrd
except ImportError:
    xlrd = None

try:
    import tabulate
except ImportError:
    tabulate = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

logger = logging.getLogger(__name__)


class DocumentIngestionDependencyError(RuntimeError):
    def __init__(self, *, filename: str, suffix: str, missing_dependencies: list[str]):
        self.filename = filename
        self.suffix = suffix
        self.missing_dependencies = list(missing_dependencies)
        super().__init__(f"{filename} 缺少文档解析依赖：{', '.join(self.missing_dependencies)}")

    def to_payload(self) -> dict:
        normalized_suffix = self.suffix.lstrip(".").lower() or "unknown"
        return {
            "code": "document_ingestion_dependencies_missing",
            "message": f"当前环境缺少解析 {self.filename} 所需的文档入库依赖。",
            "details": {
                "filename": self.filename,
                "fileType": normalized_suffix,
                "missingDependencies": self.missing_dependencies,
                "featurePackId": "document_ingestion",
                "requiredBundle": "document-ingestion",
                "recommendedNextAction": (
                    "打开管理台顶部的能力包面板，安装“文档读取能力包”，重启 V8OS 后重试；"
                    "不要调用系统 pip 安装依赖，系统 Python 不属于 V8OS 的受管运行时。"
                ),
            },
        }


class UnsupportedLegacyDocumentError(RuntimeError):
    def __init__(self, *, filename: str, suffix: str):
        self.filename = filename
        self.suffix = suffix
        super().__init__(f"{filename} 是不受支持的旧版二进制文档格式")

    def to_payload(self) -> dict:
        normalized_suffix = self.suffix.lstrip(".").lower() or "unknown"
        modern_suffix = "docx" if normalized_suffix == "doc" else "pptx"
        return {
            "code": "legacy_document_conversion_required",
            "message": f"{self.filename} 使用旧版 {normalized_suffix.upper()} 二进制格式，当前不能可靠解析。",
            "details": {
                "filename": self.filename,
                "fileType": normalized_suffix,
                "requiredFormat": modern_suffix,
                "recommendedNextAction": f"请先用受信任的 Office/LibreOffice 将文件转换为 .{modern_suffix}，再调用 read_native_file。",
            },
        }


class DocumentParser:
    """
    Unified parser for multiple document formats.
    Attempts to return content in Markdown format to leverage Semantic Chunking later.
    """
    DOCUMENT_INGESTION_DEPENDENCY_MAP = {
        ".csv": ["tabulate"],
        ".xls": ["xlrd", "tabulate"],
        ".xlsx": ["openpyxl", "tabulate"],
        ".pdf": ["PyMuPDF"],
        ".docx": ["python-docx"],
        ".pptx": ["python-pptx"],
    }

    @classmethod
    def get_missing_dependencies_for_suffix(cls, suffix: str) -> list[str]:
        normalized = str(suffix or "").lower()
        required = cls.DOCUMENT_INGESTION_DEPENDENCY_MAP.get(normalized, [])
        missing: list[str] = []
        for dependency in required:
            if dependency == "openpyxl" and openpyxl is None:
                missing.append(dependency)
            elif dependency == "xlrd" and xlrd is None:
                missing.append(dependency)
            elif dependency == "tabulate" and tabulate is None:
                missing.append(dependency)
            elif dependency == "PyMuPDF" and pymupdf is None:
                missing.append(dependency)
            elif dependency == "python-docx" and Document is None:
                missing.append(dependency)
            elif dependency == "python-pptx" and Presentation is None:
                missing.append(dependency)
        return missing

    @classmethod
    def ensure_document_ingestion_dependencies(cls, file_path: Path) -> None:
        if file_path.suffix.lower() in {".doc", ".ppt"}:
            raise UnsupportedLegacyDocumentError(
                filename=file_path.name,
                suffix=file_path.suffix,
            )
        missing = cls.get_missing_dependencies_for_suffix(file_path.suffix)
        if missing:
            raise DocumentIngestionDependencyError(
                filename=file_path.name,
                suffix=file_path.suffix,
                missing_dependencies=missing,
            )
    
    @classmethod
    def parse_file(cls, file_path: Path) -> str:
        """Parse a single file and return its textual/markdown content."""
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        cls.ensure_document_ingestion_dependencies(file_path)
            
        ext = file_path.suffix.lower()
        
        try:
            if ext in ['.txt', '.md', '.mdx']:
                return cls._parse_text(file_path)
            elif ext in ['.csv']:
                return cls._parse_csv(file_path)
            elif ext in ['.xls', '.xlsx']:
                return cls._parse_excel(file_path)
            elif ext == '.pdf':
                return cls._parse_pdf(file_path)
            elif ext == '.docx':
                return cls._parse_docx(file_path)
            elif ext == '.pptx':
                return cls._parse_pptx(file_path)
            elif ext in ['.html', '.htm']:
                return cls._parse_html(file_path)
            elif ext == '.json':
                return cls._parse_json(file_path)
            elif ext == '.xml':
                return cls._parse_text(file_path)
            elif ext in ['.zip', '.tar', '.gz']:
                return cls._parse_archive(file_path)
            else:
                logger.warning(f"[DocumentParser] Unsupported extension '{ext}', falling back to plain text.")
                return cls._parse_text(file_path)
        except (DocumentIngestionDependencyError, UnsupportedLegacyDocumentError):
            raise
        except Exception as e:
            logger.error(f"[DocumentParser] Error parsing {file_path}: {e}")
            return f"Error parsing document {file_path.name}: {e}"

    @staticmethod
    def _parse_text(file_path: Path) -> str:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    @staticmethod
    def _parse_json(file_path: Path) -> str:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            data = json.load(f)
            return json.dumps(data, indent=2, ensure_ascii=False)

    @staticmethod
    def _parse_csv(file_path: Path) -> str:
        with open(file_path, 'r', encoding='utf-8-sig', errors='ignore', newline='') as f:
            return DocumentParser._render_tabular_rows(list(csv.reader(f)))

    @staticmethod
    def _parse_excel(file_path: Path) -> str:
        if file_path.suffix.lower() == ".xlsx":
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            try:
                rows = [list(row) for row in workbook.active.iter_rows(values_only=True)]
            finally:
                workbook.close()
            return DocumentParser._render_tabular_rows(rows)

        workbook = xlrd.open_workbook(file_path, on_demand=True)
        try:
            worksheet = workbook.sheet_by_index(0)
            rows = [worksheet.row_values(index) for index in range(worksheet.nrows)]
        finally:
            workbook.release_resources()
        return DocumentParser._render_tabular_rows(rows)

    @staticmethod
    def _render_tabular_rows(rows: list[list[object]]) -> str:
        if not rows:
            return ""
        headers = ["" if value is None else value for value in rows[0]]
        body = [["" if value is None else value for value in row] for row in rows[1:]]
        return tabulate.tabulate(body, headers=headers, tablefmt="pipe", disable_numparse=True)

    @staticmethod
    def _parse_pdf(file_path: Path) -> str:
        if not pymupdf:
            return f"[PyMuPDF not installed for PDF parsing of {file_path.name}]"
        
        doc = pymupdf.open(file_path)
        text_blocks = []
        for page in doc:
            text_blocks.append(page.get_text())
        return "\n\n".join(text_blocks)

    @staticmethod
    def _parse_docx(file_path: Path) -> str:
        if not Document:
            return f"[python-docx not installed for DOCX parsing of {file_path.name}]"
        
        doc = Document(file_path)
        text_blocks: list[str] = []
        for block in doc.iter_inner_content():
            if hasattr(block, "rows"):
                rows = []
                for row in block.rows:
                    rows.append([
                        "\n".join(
                            paragraph.text.strip()
                            for paragraph in cell.paragraphs
                            if paragraph.text.strip()
                        )
                        for cell in row.cells
                    ])
                if rows:
                    width = max(len(row) for row in rows)
                    normalized_rows = [row + ([""] * (width - len(row))) for row in rows]
                    escaped_rows = [
                        [value.replace("|", "\\|").replace("\n", "<br>") for value in row]
                        for row in normalized_rows
                    ]
                    text_blocks.append("\n".join([
                        f"| {' | '.join(escaped_rows[0])} |",
                        f"| {' | '.join(['---'] * width)} |",
                        *[f"| {' | '.join(row)} |" for row in escaped_rows[1:]],
                    ]))
                continue

            style_name = str(getattr(getattr(block, "style", None), "name", "") or "")
            if style_name.startswith("Heading"):
                level = style_name.replace("Heading ", "")
                prefix = "#" * int(level) if level.isdigit() else "#"
                text_blocks.append(f"{prefix} {block.text}")
            else:
                text_blocks.append(block.text)
        return "\n\n".join(text_blocks)

    @staticmethod
    def _parse_pptx(file_path: Path) -> str:
        if not Presentation:
            return f"[python-pptx not installed for PPTX parsing of {file_path.name}]"
            
        prs = Presentation(file_path)
        text_blocks = []
        for i, slide in enumerate(prs.slides):
            text_blocks.append(f"## Slide {i+1}")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_blocks.append(shape.text)
        return "\n\n".join(text_blocks)

    @staticmethod
    def _parse_html(file_path: Path) -> str:
        if not BeautifulSoup:
            return DocumentParser._parse_text(file_path)
            
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            soup = BeautifulSoup(f.read(), "html.parser")
            
            # Simple conversion to markdown headers
            for i in range(1, 7):
                for tag in soup.find_all(f'h{i}'):
                    tag.replace_with(f"{'#' * i} {tag.get_text().strip()}\n")
            
            return soup.get_text(separator="\n\n").strip()

    @classmethod
    def _parse_archive(cls, file_path: Path) -> str:
        combined_text = []
        with tempfile.TemporaryDirectory() as temp_dir:
            if file_path.suffix.lower() == '.zip':
                with zipfile.ZipFile(file_path, 'r') as zip_ref:
                    zip_ref.extractall(temp_dir)
            elif file_path.suffix.lower() in ['.tar', '.gz']:
                with tarfile.open(file_path, 'r:*') as tar_ref:
                    def is_within_directory(directory, target):
                        abs_directory = os.path.abspath(directory)
                        abs_target = os.path.abspath(target)
                        prefix = os.path.commonprefix([abs_directory, abs_target])
                        return prefix == abs_directory

                    def safe_extract(tar, path=".", members=None, *, numeric_owner=False):
                        for member in tar.getmembers():
                            member_path = os.path.join(path, member.name)
                            if not is_within_directory(path, member_path):
                                raise Exception("Attempted Path Traversal in Tar File")
                        tar.extractall(path, members, numeric_owner=numeric_owner)
                    
                    safe_extract(tar_ref, path=temp_dir)
            
            for root, dirs, files in os.walk(temp_dir):
                for file in files:
                    extracted_path = Path(root) / file
                    # Only parse supported extensions to avoid binary garbage
                    if extracted_path.suffix.lower() in ['.txt', '.md', '.mdx', '.csv', '.pdf', '.doc', '.docx', '.ppt', '.pptx', '.html', '.htm', '.json', '.xml']:
                        content = cls.parse_file(extracted_path)
                        combined_text.append(f"# File: {file}\n\n{content}")
                        
        return "\n\n---\n\n".join(combined_text)

document_parser = DocumentParser()
